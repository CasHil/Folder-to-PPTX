import os
import argparse
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

def add_centered_image_slide(presentation, image_path):
    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)
    
    img = Image.open(image_path)
    img_width_px, img_height_px = img.size
    img_aspect_ratio = img_width_px / img_height_px
    
    slide_width_inch = presentation.slide_width.inches
    slide_height_inch = presentation.slide_height.inches
    
    if img_width_px > img_height_px:
        img_width_inch = slide_width_inch * 0.9
        img_height_inch = img_width_inch / img_aspect_ratio
        if img_height_inch > slide_height_inch * 0.9:
            img_height_inch = slide_height_inch * 0.9
            img_width_inch = img_height_inch * img_aspect_ratio
    else:
        img_height_inch = slide_height_inch * 0.9
        img_width_inch = img_height_inch * img_aspect_ratio
        if img_width_inch > slide_width_inch * 0.9:
            img_width_inch = slide_width_inch * 0.9
            img_height_inch = img_width_inch / img_aspect_ratio
    
    left_inch = (slide_width_inch - img_width_inch) / 2
    top_inch = (slide_height_inch - img_height_inch) / 2

    slide.shapes.add_picture(image_path, Inches(left_inch), Inches(top_inch), Inches(img_width_inch), Inches(img_height_inch))

def create_centered_images_presentation(directory, output_file):
    presentation = Presentation()
    
    for root, dirs, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.wmf')):
                image_path = os.path.join(root, file)
                add_centered_image_slide(presentation, image_path)
    
    presentation.save(output_file)
    print(f"Presentation saved as {output_file} with {len(presentation.slides)} slides.")

def main():
    parser = argparse.ArgumentParser(description="Create a PowerPoint presentation from images in a directory.")
    parser.add_argument('--directory', type=str, default=os.getcwd(), help="The directory containing images.")
    parser.add_argument('--filename', type=str, default='presentation.pptx', help="The filename for the output presentation.")
    args = parser.parse_args()

    create_centered_images_presentation(args.directory, args.filename)

if __name__ == "__main__":
    main()
